from flask import Flask, request, jsonify, render_template, redirect, url_for, send_from_directory
import os
import time
import multiprocessing
from concurrent.futures import ThreadPoolExecutor, as_completed
import openai
import requests
import json
import pytesseract
from pdf2image import convert_from_path
from PyPDF2 import PdfReader
import re
import logging
from pptx import Presentation

app = Flask(__name__, template_folder='templates')

# OpenAI API-Schlüssel festlegen
api_key = "My_API_TOKEN"
openai.api_key = api_key

logger = logging.getLogger(__name__)

# Setze die Pfade
tesseract_path = "/app/Tesseract/tesseract.exe"
pytesseract.pytesseract.tesseract_cmd = tesseract_path

poppler_path = "/app/Poppler/bin"

# Überprüfe, ob Tesseract existiert
if not tesseract_path or not os.path.isfile(tesseract_path):
    logger.error(f"Tesseract-Pfad ist nicht korrekt gesetzt oder die Datei existiert nicht. Pfad: {tesseract_path}")
    print(f"Error: Tesseract-Pfad ist nicht korrekt gesetzt oder Datei fehlt. Pfad: {tesseract_path}")
    exit(1)
else:
    logger.info(f"Tesseract-Pfad korrekt gesetzt: {tesseract_path}")

# Überprüfe, ob Poppler existiert
if not poppler_path or not os.path.isdir(poppler_path):
    logger.error(f"Poppler-Pfad ist nicht korrekt gesetzt oder das Verzeichnis existiert nicht. Pfad: {poppler_path}")
    print(f"Error: Poppler-Pfad ist nicht korrekt gesetzt oder Verzeichnis fehlt. Pfad: {poppler_path}")
    exit(1)
else:
    logger.info(f"Poppler-Pfad korrekt gesetzt: {poppler_path}")

# Set the folder to save uploaded files
UPLOAD_FOLDER = os.path.join(os.getcwd(), 'uploads')
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

def get_pdf_page_count(pdf_path):
    with open(pdf_path, 'rb') as file:
        reader = PdfReader(file)
        return len(reader.pages)

def extract_text_from_page(pdf_path, page_index):
    try:
        # Stelle sicher, dass poppler_path korrekt gesetzt ist
        images = convert_from_path(pdf_path, first_page=page_index + 1, last_page=page_index + 1, poppler_path=poppler_path, dpi=150)
        if images:
            # OCR mit Tesseract durchführen
            page_text = pytesseract.image_to_string(images[0], lang='deu')
            return page_index, page_text
        else:
            logger.warning(f"No images found on page {page_index + 1}")
            return page_index, ""
    except Exception as e:
        # Detaillierte Fehlerprotokollierung
        logger.error(f"Fehler bei der Textextraktion auf Seite {page_index + 1}: {e}", exc_info=True)
        return page_index, ""

def extract_text_from_pdf(pdf_path):
    total_pages = get_pdf_page_count(pdf_path)
    text = [""] * total_pages    
    with ThreadPoolExecutor(max(total_pages, multiprocessing.cpu_count())) as executor:
        futures = {executor.submit(extract_text_from_page, pdf_path, i): i for i in range(total_pages)}
        for future in as_completed(futures):
            page_index, page_text = future.result()
            text[page_index] = page_text

    full_text = "\n\n".join([text[i] for i in range(total_pages)])
    return full_text

def create_detailed_summary(api_key, text, temperature=0.6):
    url = "https://api.openai.com/v1/chat/completions"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}"
    }
    summary_prompt = f"""
Fasse den folgenden Text ausführlich zusammen und halte dich dabei strikt an das vorgegebene standardisierte Format. Die Zusammenfassung soll in vollständigen Sätzen erfolgen, du hast dabei folgende Anweisungen:
Umfassende Berücksichtigung: Erfasse alle relevanten Informationen und Details des Textes, ohne dass wesentliche Inhalte verloren gehen.
Formatvorgabe einhalten: Nutze ausschließlich das vorgegebene standardisierte Format. Füge keine zusätzlichen Stichpunkte, Anführungszeichen oder andere Elemente ein, die das Format oder die Strukturierung beeinträchtigen.
Eingabefelder korrekt nutzen: Trage die relevanten Informationen in die vorgesehenen Felder ein, die durch [ ] markiert sind.
Wichtige Hinweise: Gib keine IBANs aus, erwähne keine spezifischen Namen von AGBs, und verwende keine Anführungszeichen. Achte darauf, alle Sicherheiten vollständig aufzulisten.
Kreditart: Gib bei [Art des Kredites] ausschließlich die Kreditgattung an, zum Beispiel "Darlehen" oder "Festdarlehen", je nachdem, was im Vertrag steht.
Sicherheiten: Schreibe bei [Sicherheiten] alle Sicherheiten in einen Satz, getrennt durch Kommata
Formulierungen: Gib Zahlen immer im Tausenderformat an. Also bspw. "3,4T €" statt 34.000€
Ziel: Erstelle eine präzise und strukturierte Zusammenfassung, die den Text vollständig und korrekt wiedergibt, ohne das standardisierte Format zu verlassen.

Hier ist der zu zusammenfassende Text:
    
    {text}

    Standardisiertes Format:
    • Art des Kredits: [Art des Kredits] 
    • Kreditinstitut: [Kreditinstitut] 
    • Kunde: [Kunde] 
    • Datum: [Datum]
    • Kreditrahmen: [Betrag]
    • Sollzinssatz: [Sollzinssatz]
    • Nebenleistungen: [Nebenleistungen]
    • Darlehensrückzahlung und Laufzeit: [Darlehensrückzahlung und Laufzeit]
    • Sicherheiten: [Sicherheiten]
    • Sonstiges: [Sonstiges]
    • AGB/Covenants/Bemerkungen: [Bemerkungen]
    """
    data = {
        "model": "gpt-4o-mini",
        "messages": [
            {"role": "system", "content": "You are a helpful assistant who creates detailed and standardized summaries of loan agreements."},
            {"role": "user", "content": summary_prompt}
        ],
        "max_tokens": 750,
        "temperature": temperature,
        "n": 1
    }

    response = requests.post(url, headers=headers, data=json.dumps(data))
    if response.status_code == 200:
        response_data = response.json()
        tokens_used = response_data["usage"]["total_tokens"]
        summary = response_data["choices"][0]["message"]["content"]
        return summary, tokens_used
    else:
        logger.error(f"Ein Fehler ist aufgetreten: {response.status_code} - {response.text}")
        return "", 0

def validate_and_correct_summaries(api_key, summary1, summary2, original_text):
    url = "https://api.openai.com/v1/chat/completions"
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {api_key}"
    }
    validation_prompt = f"""
Anweisung:
Erstelle eine konsolidierte, detaillierte Zusammenfassung aus den folgenden zwei Texten, in der sämtliche Informationen aus beiden Texten enthalten sind. Achte darauf, dass keine Details verloren gehen und alle relevanten Informationen erhalten bleiben.

Vorgehensweise:

Informationsintegration: Integriere alle Informationen aus beiden Zusammenfassungen in eine einzige, vollständige Zusammenfassung. Es dürfen keine Details ausgelassen werden.
Konfliktlösung: Bei widersprüchlichen Angaben vergleiche die Informationen mit dem Originaltext, um die korrekte Information zu ermitteln und in die Zusammenfassung einzufügen.
Sprachliche Präzision: Achte auf korrekte Rechtschreibung, Grammatik und Syntax.
Strukturierte Ausgabe: Nutze das vorgegebene standardisierte Format für die Ausgabe (siehe unten) und stelle sicher, dass der gesamte Inhalt korrekt verstanden und erfasst wurde.
Formulierungen: Gib Zahlen immer im Tausenderformat an. Also bspw. "3,4T €" statt 34.000€

[Sicherheiten]: Notiere hier alle genannten Sicherheiten.
[Art des Kredits]: Gib hier nur die Kreditart (z.B. „Darlehen“) an, ohne zusätzliche Details.
[Kreditinstitut]: Führe hier nur den Namen des Kreditinstituts auf.
[Kunde]: Gib hier nur den Namen des Kunden an.
[Datum des Vertrages]: Führe hier das Datum des Vertragsabschlusses auf.
[Sollzinssatz]: Notiere den Sollzinssatz sowie alle vorhandenen variablen Komponenten.
[Bemerkungen]: Erstelle eine vollständige Liste aller Covenants, AGB und sonstiger Bemerkungen aus dem Vertrag.
[Sonstiges]: Erfasste hier alle relevanten Informationen, die nicht unter den anderen Kategorien abgedeckt sind.

Erstelle eine präzise und umfassende Zusammenfassung, die alle wichtigen Details enthält, korrekt strukturiert ist und den Anforderungen des vorgegebenen Formats entspricht.

    Zusammenfassung 1:
    {summary1}

    Zusammenfassung 2:
    {summary2}

    Originaltext:
    {original_text}

    Standardisiertes Format:
    • Art des Kredits: [Art des Kredits] 
    • Kreditinstitut: [Kreditinstitut] 
    • Kunde: [Kunde] 
    • Datum: [Datum des Vertrages]
    • Kreditrahmen: [Betrag]
    • Sollzinssatz: [Sollzinssatz]
    • Nebenleistungen: [Nebenleistungen]
    • Darlehensrückzahlung und Laufzeit: [Darlehensrückzahlung und Laufzeit]
    • Sicherheiten: [Sicherheiten]
    • Sonstiges: [Sonstiges]
    • AGB/Covenants/Bemerkungen: [Bemerkungen]
    """
    data = {
        "model": "gpt-4o-mini",
        "messages": [
            {"role": "system", "content": "Du bist ein hilfreicher Assistent, der Zusammenfassungen von Kreditverträgen konsolidiert. Achte besonders auf Details und behalte alle Informationen bei. Antworte immer auf Deutsch."},
            {"role": "user", "content": validation_prompt}
        ],
        "max_tokens": 1000,
        "temperature": 0.3,
        "n": 1
    }

    response = requests.post(url, headers=headers, data=json.dumps(data))
    if response.status_code == 200:
        response_data = response.json()
        tokens_used = response_data["usage"]["total_tokens"]
        consolidated_summary = response_data["choices"][0]["message"]["content"]
        return consolidated_summary, tokens_used
    else:
        logger.error(f"Ein Fehler ist aufgetreten: {response.status_code} - {response.text}")
        return "", 0

def extract_information(content, index):
    patterns = {
        'Art des Kredits:': f'UXHGA01{index}',
        'Kreditinstitut:': f'UXHGA02{index}',
        'Kunde:': f'UXHGA03{index}',
        'Datum:': f'UXHGA04{index}',
        'Kreditrahmen:': f'UXHGD01{index}',
        'Sollzinssatz:': f'UXHGE01{index}',
        'Nebenleistungen:': f'UXHGF01{index}',
        'Darlehensrückzahlung und Laufzeit:': f'UXHGG01{index}',
        'Sicherheiten:': f'UXHGH01{index}',
        'Sonstiges:': f'UXHGI01{index}',
        'AGB/Covenants/Bemerkungen:': f'UXHGJ01{index}'
    }

    extracted_data = {}

    for label, placeholder in patterns.items():
        pattern = re.escape(label) + r'\s*(.*?)(?:\n|•|$)'
        match = re.search(pattern, content)
        if match:
            extracted_value = match.group(1).strip()
            extracted_value = re.sub(r'^[-*•\s]+', '', extracted_value)  # Entfernt führende Sonderzeichen und Leerzeichen
            extracted_data[placeholder] = extracted_value
            logger.info(f"Successfully extracted: {placeholder} = {extracted_data[placeholder]}")
        else:
            logger.warning(f"Could not extract {placeholder}. Pattern not found.")

    return extracted_data

def replace_text_in_shape(shape, extracted_data):
    if not hasattr(shape, "text_frame"):
        return False
    text_frame = shape.text_frame
    modified = False
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            original_text = run.text
            for key, value in extracted_data.items():
                if key in original_text:
                    new_text = original_text.replace(key, value)
                    if new_text != original_text:
                        run.text = new_text
                        logger.info(f"Replaced '{key}' with '{value}' in shape text.")
                        modified = True
    return modified

def update_presentation(prs, extracted_data_list):
    for extracted_data in extracted_data_list:
        for slide in prs.slides:
            if slide.shapes:
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                replace_text_in_shape(cell, extracted_data)
                    else:
                        replace_text_in_shape(shape, extracted_data)

def calculate_costs(input_tokens, output_tokens):
    cost_input = (input_tokens / 1_000_000) * 0.15
    cost_output = (output_tokens / 1_000_000) * 0.6
    return cost_input + cost_output

def process_single_document(pdf_path, index):
    logger.info(f"*Verarbeitung von Vertrag {index} beginnt...*")

    # Schritt 1: Extrahieren des Textes aus dem Originalvertrag
    logger.info(f"Schritt 1: Inhalte des {index}. Vertrags werden gelesen...")
    original_text = extract_text_from_pdf(pdf_path)
    logger.info(f"Originaltext des {index}. Vertrags:\n{original_text}")

    # Schritt 2: Separate Analyse des Originalvertrags
    logger.info(f"Schritt 2: Erste Analyse des Originalvertrags {index} wird durchgeführt...")
    summary1, tokens_used1 = create_detailed_summary(api_key, original_text, 0.4)
    logger.info(f"Erste Zusammenfassung des {index}. Vertrags:\n{summary1}")

    logger.info(f"Schritt 2: Zweite Analyse des Originalvertrags {index} wird durchgeführt...")
    summary2, tokens_used2 = create_detailed_summary(api_key, original_text, 0.3)
    logger.info(f"Zweite Zusammenfassung des {index}. Vertrags:\n{summary2}")

    # Schritt 4: Konsolidierung der Analysen des Originalvertrags
    logger.info(f"Schritt 4: Konsolidierung der beiden Analysen des Originalvertrags {index}...")
    consolidated_original_summary, tokens_used_consolidation = validate_and_correct_summaries(api_key, summary1, summary2, original_text)
    logger.info(f"Konsolidierte Zusammenfassung des Originalvertrags {index}:\n{consolidated_original_summary}")

    # Schritt 6: Extrahieren der relevanten Informationen aus der endgültigen Zusammenfassung
    logger.info(f"Schritt 6: Extrahieren der relevanten Informationen für Vertrag {index} ...")
    extracted_data = extract_information(consolidated_original_summary, index)

    logger.info(f"Finalisierung der Vertragsanalyse für Vertrag {index} abgeschlossen.")
    return extracted_data

def process_multiple_documents(document_sets, pptx_path):
    extracted_data_list = []
    contracts_count = len(document_sets)

    logger.info(f"*Anzahl der zu verarbeitenden Verträge: {contracts_count}*")  # Logging der Anzahl der Verträge

    logger.info("*Parallelisierte Verarbeitung aller Verträge wird gestartet...*")

    # Sicherstellen, dass die Reihenfolge beibehalten wird
    with ThreadPoolExecutor(max_workers=min(len(document_sets), multiprocessing.cpu_count())) as executor:
        futures = [
            executor.submit(process_single_document, document_set['pdf_path'], index + 1)
            for index, document_set in enumerate(document_sets)
        ]
        for future in as_completed(futures):
            extracted_data_list.append(future.result())

    logger.info("*Aktualisieren der PowerPoint-Präsentation...*")
    prs = Presentation(pptx_path)
    update_presentation(prs, extracted_data_list)

    # Änderungen direkt in der ursprünglichen Datei speichern
    prs.save(pptx_path)
    logger.info(f"Präsentation aktualisiert und gespeichert: {pptx_path}")

    return pptx_path  # Return the path to the updated PowerPoint file

@app.route('/upload', methods=['POST'])
def upload_files():
    if 'files[]' not in request.files or 'pptFile' not in request.files:
        return redirect(request.url)

    ppt_file = request.files['pptFile']

    # Verwenden einer Menge, um doppelte Pfade zu vermeiden
    file_paths_set = set()

    document_sets = []
    for i, file in enumerate(request.files.getlist('files[]')):
        file_path = os.path.join(UPLOAD_FOLDER, file.filename)
        
        # Überprüfen, ob der Dateipfad bereits existiert
        if file_path not in file_paths_set:
            file.save(file_path)
            logger.info(f"Vertrag {i+1}: {file.filename} erfolgreich hochgeladen und gespeichert.")
            file_paths_set.add(file_path)
            
            document_sets.append({
                'pdf_path': file_path,
                'addenda_paths': []  # Nachträge werden ignoriert
            })
        else:
            logger.warning(f"Datei {file.filename} wurde bereits hochgeladen, wird ignoriert.")

    logger.info(f"Anzahl der hochgeladenen Verträge: {len(document_sets)}")

    ppt_file_path = os.path.join(UPLOAD_FOLDER, ppt_file.filename)
    ppt_file.save(ppt_file_path)

    # Generiere eine eindeutige Task-ID
    task_id = str(time.time())

    # Verarbeiten der Dateien in einem separaten Thread
    with ThreadPoolExecutor(max_workers=1) as executor:
        future = executor.submit(process_multiple_documents, document_sets, ppt_file_path)
        processed_pptx_path = future.result()

    # Return the download link and the task_id for progress tracking
    return jsonify({
        'message': 'Processing complete',
        'download_link': url_for('download', filename=os.path.basename(processed_pptx_path))
    })
    
@app.route('/')
def index():
    return render_template('index.html')

@app.route('/download/<filename>', methods=['GET'])
def download(filename):
    try:
        logger.info(f"Attempting to send file from {os.path.join(UPLOAD_FOLDER, filename)}")
        return send_from_directory(directory=UPLOAD_FOLDER, path=filename, as_attachment=True)
    except Exception as e:
        logger.error(f"Error sending file: {e}")
        return "File not found", 404

if __name__ == '__main__':
    if not os.path.exists(UPLOAD_FOLDER):
        os.makedirs(UPLOAD_FOLDER)
    app.run(debug=True)
